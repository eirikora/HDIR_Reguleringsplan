from selenium import webdriver
from selenium.webdriver.chrome.options import Options
import csv
import re
import os
import time
from pathlib import Path
import mimetypes
import traceback
from bs4 import BeautifulSoup
import subprocess
from urllib.parse import urlparse, parse_qs, unquote, quote, urljoin
import mammoth
import openpyxl
from pptx import Presentation
import shutil

# Read the CSV file and extract URLs
csv_file = 'download-regulation-reports.csv'

# A filter variable to download only a specific url in csv file. Set to None if all files in the csv file
only_one_url = None
only_one_url = 'https://www.ehelse.no/kodeverk-og-terminologi/Norsk-klinisk-prosedyrekodeverk-(NKPK)'
# only_one_url = 'https://finnkode.helsedirektoratet.no/adm/collections/3512'
# only_one_url = 'https://health.ec.europa.eu/document/download/b744f30b-a05e-4b9c-9630-ad96ebd0b2f0_en?filename=ehn_guidelines_eprescriptions_en.pdf'

# Name of database of downloaded files
filedatabase = 'downloaded_files.csv'

# Set to keep track of downloaded files
downloaded_file_urls = {}

# Set up Chrome options
chrome_options = Options()
chrome_options.add_argument("--headless")
chrome_options.add_argument("--disable-gpu")
chrome_options.add_argument("--no-sandbox")
chrome_options.add_argument("--disable-dev-shm-usage")
chrome_options.add_argument(
    "user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/115.0 Safari/537.36"
)

# Initialize the WebDriver
driver = webdriver.Chrome(options=chrome_options)

with open(csv_file, 'r', encoding='iso-8859-1') as file:
    reader = csv.DictReader(file)
    url_dict = {}
    for row in reader:
        url = row.get('referanse_url')
        if url and url not in url_dict:
            referanse_lenketekst = row.get('referanse_lenketekst', '')
            url_dict[url] = referanse_lenketekst

# Function to clean filenames
def clean_filename(filename, max_length=100):
    filename, _ = os.path.splitext(filename)
    filename = re.sub(r'\s*\(.*?\)\s*', '', filename)
    filename = re.sub(r'[<>:"/\\|?*]', '', filename)
    filename = filename.strip().replace(' ', '_')
    if len(filename) > max_length:
        filename = filename[:max_length]
    return filename

# Function to download a file using curl
def download_with_curl(url, filepath):
    try:
        command = [
            'curl',
            '-L',
            '-o', filepath,
            '-H', 'User-Agent: Mozilla/5.0 (Windows NT 10.0; Win64; x64)',
            url
        ]
        result = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if result.returncode == 0:
            print(f'Downloaded with curl: {os.path.basename(filepath)}')
            return True
        else:
            print(f'Curl failed to download {url}: {result.stderr.decode()}')
            return False
    except Exception as e:
        print(f'Exception during curl download of {url}: {e}')
        return False

# Function to determine if URL is an HTML page
def is_html_url(url):
    mime_type, _ = mimetypes.guess_type(url)
    if mime_type:
        return mime_type.startswith('text/html')
    else:
        other_extensions = ['.txt', '.pdf', '.xml', '.xsd', '.xls', '.doc', '.ppt', '.xlsx', '.docx', '.pptx']
        path = urlparse(url).path
        ext = os.path.splitext(path)[1].lower()
        return not ext in other_extensions

# Function to convert .docx to HTML using Mammoth
def convert_docx_to_html(filepath):
    try:
        with open(filepath, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html = result.value
            html_filename = os.path.splitext(filepath)[0] + ".html"
            with open(html_filename, "w", encoding="utf-8") as html_file:
                html_file.write(html)
            print(f"Converted {filepath} to HTML: {html_filename}")
        print('Conversion completed. Removing original docx file')
        # Check if file is closed before removing
        if not docx_file.closed:
            docx_file.close()
        os.remove(filepath)
        return html_filename
    except Exception as e:
        print(f"Failed to convert {filepath} to HTML: {e}")
        return filepath

# Function to convert .xlsx to .txt using openpyxl
def convert_xlsx_to_txt(filepath):
    try:
        wb = openpyxl.load_workbook(filepath)
        txt_filename = os.path.splitext(filepath)[0] + ".txt"
        with open(txt_filename, "w", encoding="utf-8") as txt_file:
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                txt_file.write(f"Sheet: {sheet}\n")
                for row in ws.iter_rows(values_only=True):
                    txt_file.write("\t".join([str(cell) if cell else "" for cell in row]) + "\n")
        print(f"Converted {filepath} to TXT: {txt_filename}")
        os.remove(filepath)
        return txt_filename
    except Exception as e:
        print(f"Failed to convert {filepath} to TXT: {e}")
        return filepath

# Function to convert .pptx to .txt using python-pptx
def convert_pptx_to_txt(filepath):
    try:
        prs = Presentation(filepath)
        txt_filename = os.path.splitext(filepath)[0] + ".txt"
        with open(txt_filename, "w", encoding="utf-8") as txt_file:
            for slide_num, slide in enumerate(prs.slides, start=1):
                txt_file.write(f"Slide {slide_num}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        txt_file.write(shape.text + "\n")
                txt_file.write("\n")
        print(f"Converted {filepath} to TXT: {txt_filename}")
        os.remove(filepath)
        return txt_filename
    except Exception as e:
        print(f"Failed to convert {filepath} to TXT: {e}")
        return filepath

# Function to convert to format supported by OpenAI
def convert_to_supported_format(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    # Handle conversion based on file type
    if ext == '.docx':
        filepath = convert_docx_to_html(filepath)
    elif ext == '.xlsx':
        filepath = convert_xlsx_to_txt(filepath)
    elif ext == '.pptx':
        filepath = convert_pptx_to_txt(filepath)
    elif ext == '.xsd':
        txt_filepath = filepath + '.txt'
        shutil.move(filepath, txt_filepath)
        filepath = txt_filepath
        print(f'Renamed file to: {txt_filepath}')
    # Returning new filepath as file may have changed as it was converted
    return filepath

# Read CSV database of existing files
if Path(filedatabase).exists():
    print("\nReading the existing database of downloaded files:")
    with open(filedatabase, 'r', newline='', encoding='utf-8') as csvfile:
        reader = csv.DictReader(csvfile)
        for row in reader:
            downloaded_file_urls[row['Filename']] = row['URL']
    csvfile.close()
    print(f"Read {len(downloaded_file_urls)} files from database.")

# Create output directory
output_dir = 'nedlastede_filer'
os.makedirs(output_dir, exist_ok=True)

# List of file extensions to download
download_file_extensions = ['.pdf', '.docx', '.pptx', '.xlsx', '.xsd']

# Initialize sequence number
sequence_number = 0

# Maximum filename length
max_filename_length = 80

# Download and save files
for url, ref_text in url_dict.items():
    sequence_number += 1
    if not only_one_url or url == only_one_url:
        try:
            seq_str = f"{sequence_number:03d}_"
            print(f'Processing {url}')

            if is_html_url(url):
                driver.get(url)
                time.sleep(10)
                rendered_html = driver.execute_script("return document.documentElement.outerHTML;")

                # Strip down the HTML as much as possible by removing elements of no need for AI/Vector store, like scripts, buttons, style, etc.
                soup = BeautifulSoup(rendered_html, 'html.parser')

                for tag in soup(['script', 'style', 'header', 'footer', 'nav', 'aside', 'img', 'button', 'input']):
                    tag.decompose()

                # Strip all content in tags above, except for href attributes in a-tags
                for tag in soup.find_all():
                    if tag.name == 'a':
                        tag.attrs = {key: value for key, value in tag.attrs.items() if key == 'href'}
                    else:
                        tag.attrs = {}

                parsed_url = urlparse(url)
                base_url = f"{parsed_url.scheme}://{parsed_url.netloc}"

                # Make all hyperlinks non-local and consider downloading if interesting sub-documents 
                for a_tag in soup.find_all('a', href=True):
                    href = a_tag['href']
                    absolute_href = urljoin(base_url, href)
                    a_tag['href'] = absolute_href

                    parsed_href = urlparse(absolute_href)
                    ext = os.path.splitext(parsed_href.path)[1].lower()

                    # Define a suitable filename for the subfile pointed to and to be downloaded
                    filename_list = parse_qs(parsed_href.query).get('filename')
                    if not filename_list:
                        thefilename = os.path.basename(parsed_href.path)
                    else:
                        thefilename = filename_list[0]
                    thefilename = unquote(thefilename)

                    absolute_href = unquote(absolute_href)
                    encoded_abs_url = quote(absolute_href, safe=":/")
                    
                    # If the hypertext link points to a sub-file that we have not seen yet, download it too
                    if ext in download_file_extensions and encoded_abs_url not in downloaded_file_urls.values():
                        filename = clean_filename(thefilename, max_length=max_filename_length) + ext
                        filename = seq_str + filename
                        filepath = os.path.join(output_dir, filename)
                        success = download_with_curl(encoded_abs_url, filepath)
                        if success:
                            # If successful download, convert the file to a format supported by OpenAI Vector Stores
                            filepath = convert_to_supported_format(filepath)
                            downloaded_file_urls[filepath] = encoded_abs_url
                        else:
                            print(f'Failed to download linked file: {encoded_abs_url}')

                # Convert beautiful soup structure to cleaned text
                main_content = soup.find('div', {'id': 'root'})
                if main_content:
                    clean_html = str(main_content)
                else:
                    clean_html = str(soup)

                # Save html file
                filename = seq_str + clean_filename(ref_text, max_length=max_filename_length) + '.html'
                filepath = os.path.join(output_dir, filename)
                with open(filepath, 'w', encoding='utf-8') as file:
                    file.write(clean_html)
                downloaded_file_urls[filepath] = url
                print(f'Saved {filename}')
            else:
                # Download a non-html URL with CURL
                parsed_url = urlparse(url)
                path = unquote(parsed_url.path)
                ext = os.path.splitext(url)[1]

                filename = seq_str + clean_filename(ref_text, max_length=max_filename_length) + ext
                filepath = os.path.join(output_dir, filename)

                success = download_with_curl(url, filepath)
                if success:
                    print(f'Downloaded file: {filename}')
                    filepath = convert_to_supported_format(filepath)
                    downloaded_file_urls[filepath] = url
                else:
                    print(f'Failed to download file: {url}')

        except Exception as e:
            print(f'Could not process {url}: {e}')
            traceback.print_exc()

driver.quit()


# Writing to CSV using DictWriter
with open(filedatabase, 'w', newline='', encoding='utf-8') as csvfile:
    fieldnames = ['Filename', 'URL']
    writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
    
    writer.writeheader()
    for filename, url in downloaded_file_urls.items():
        writer.writerow({'Filename': filename, 'URL': url})